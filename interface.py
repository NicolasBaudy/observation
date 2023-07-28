# -*- coding: utf-8 -*-
import streamlit as st
import streamlit.components.v1 as components


import pandas as pd
import seaborn as sns
import matplotlib.pyplot as plt
import random
import numpy as np
from datetime import datetime, timedelta
from scipy.signal import savgol_filter
# from pydrive.auth import GoogleAuth
# from pydrive.drive import GoogleDrive
# from io import StringIO
import os

# import base64
# from pandas.plotting import table as pd_table
# from pdflatex import PDFLaTeX
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

import pystan
from fbprophet import Prophet


# import base64
# from pandas.plotting import table as pd_table
# from pdflatex import PDFLaTeX

# POUR LANCER L'INTERFACE EN LOCAL:
#   streamlit run interface.py

# POUR LANCER L'INTERFACE SUR LE WEB, après avoir mis le code sur le dépot 
# https://github.com/Lee-RoyMannier/tourisme
# https://share.streamlit.io/lee-roymannier/tourisme/main/interface.py

# st.set_option('deprecation.showPyplotGlobalUse', False)

# TODO: 
# Se référer à 2019 pour faire les prévisions. 

### I - LECTURE DES DONNEES 
def lecture_donnees(data):
    # Formatage de l'index en date
    data = data.set_index(data.columns[0])
    data.index = data.index.map(lambda x: datetime.strptime(x, "%Y-%m-%d").date())

    # Formatage des nombres à vigule en flottant
    data = data.applymap(lambda x: float(x.replace(",", ".")))

    return data

def ordre_alpha(categorie):
    """ Pour faciliter la navigation parmi les fichiers, ces derniers sont
    classés par ordre alphabétique. On réorganisera ainsi les paires de
    clé/valeur du dictionnaire 'categorie'."""
    ordonne = sorted(categorie.items(), key=lambda x: x[0])
    categorie = {}
    for donnee in ordonne:
        categorie[donnee[0]] = donnee[1]
    return categorie

def convertion_nom_pays(correspondances_pays, code_iso):
    """ Nom en Français d'un pays à partir de son code iso en 2 lettres.
    Retourne par exemple "France" pour "FR" """
    try:
        nom_converti = correspondances_pays.loc[code_iso]["nom_pays"]
        return nom_converti
    except: 
        return code_iso

#@st.cache(persist=False)
def acquisition_donnees():
    # Code iso des pays traduits en noms français courts à partir d'un fichier
    pays = pd.read_csv("iso-pays.csv", header=None)
    pays = pays[[2,4]]
    pays.columns = ["iso", "nom_pays"]
    pays = pays.set_index("iso")
    
    # Lecture des fichiers des tables d'analyse et de leurs noms respectifs.
    # On parcoure pour cela les dossiers de données , organisés en trois
    # principales catégories: les destinations françaises, toutes les
    # destinations et une analyse génériques.  
    data = {}
    emplacement = os.path.join("data_tourisme")
    dossiers_source = ['generiques',
                       'destinations_francaises',
                       'Toutes_destinations']
    for dossier in dossiers_source:
        data[dossier] = {}
        source = os.path.join("data_tourisme/"+dossier)
        sous_dossier = os.listdir(source)[-1]

        data_dossier = "/".join([emplacement, dossier, sous_dossier])
        for donnee_tourisme in os.listdir(data_dossier):
            try:
                donnees_brut = data_dossier + "/" + donnee_tourisme
                analyse = pd.read_csv(donnees_brut, sep=";",
                                      encoding="ISO-8859-1",
                                      engine='python')
                
                # Le nom du fichier est décomposé pour former le nom qui sera affiché
                decompose = donnee_tourisme.split("_")
                type_analyse = decompose[1]
                type_analyse = type_analyse.split("-")
                nouv_type_analyse = type_analyse[1]
                
                # Les analyses générales
                if type_analyse[0] == "Generique":
                    data[dossier][nouv_type_analyse] = analyse
                    
                # Les analyses par pays
                else:
                    nom_pays = convertion_nom_pays(pays, decompose[0])
                    if not nom_pays in data[dossier].keys():
                        data[dossier][nom_pays] = {}
                    data[dossier][nom_pays][nouv_type_analyse] = analyse
                    
            except:
                pass
           
    # Réorganisation par ordre alphabétique des données
    for type_analyse in data:
        data[type_analyse] = ordre_alpha(data[type_analyse])
        if type_analyse != "generiques":
            for pays in data[type_analyse]:
                data[type_analyse][pays] = ordre_alpha(data[type_analyse][pays])
    return data, dossiers_source


def donnees_aleatoires(t0=datetime(2017, 6, 1).date(), nb_semaines=4*53):
    data = pd.DataFrame()
 
    data.index = [t0+timedelta(days=i*7) for i in range(nb_semaines)]
    for pays in ['FR', 'SU', 'EN', 'IT', 'ES']:
        data[pays] = [random.gauss(10, 4) for i in range(nb_semaines)]

    data.index.name = "Paris"

    return data


### II - MISE EN FORME

nom_pays_modif = {
    'Guadeloupe': 'OM-Antilles',
    'Nouvelle' : 'OM-Pacifique',
    'Mayotte' : f'OM-Oc.Indien'
    }

def find_key(v): 
    result = v
    for k, val in nom_pays_modif.items(): 
        if v == val: 
            result = k 
    return result

def changement_nom(pays_nom):
    if (pays_nom in nom_pays_modif.keys()):
        return nom_pays_modif[pays_nom]
    else:
        return pays_nom
    
def changement_OM(col_df):
    i = 0
    list_col = col_df
    for nom in list_col:
        list_col[i] = changement_nom(nom)
        i = i+1
    return list_col


month_str = {
    1: "janvier" , 2: "février"  , 3: "mars", 
    4: "avril"   , 5: "mai"      , 6: "juin", 
    7: "juillet" , 8: "août"     , 9: "septembre",
    10:"octobre" , 11:"novembre" , 12:"décembre"}


def duree_str(date1, date2):
    """Ecrit l'interval entre deux dates d'une manière intuitive et sans 
    redondances. Les dates en entrée sont au format de la librairie datetime.
    
    Par exemple, si on en en entrée:
    >>> date1 = datetime(2020, 10, 3)
    >>> date2 = datetime(2020, 10, 10)
    
    Alors l'interval entre les deux date s'écrira: 'du 3 au 10 octobre 2020' à
    la place par exemple de l'écriture redondante: 'du 3 octobre 2020 au 10 
    octobre 2020'.
    
    Si cela est nécessaire, les années et les mois sont précisés pour chaque
    date. Par exemple, on écrira: 'du 3 octobre 2020 au 10 septembre 2021'."""
    
    d1 = min(date1, date2)
    d2 = max(date1, date2)
    
    def day_str(j):
        if j==1:
            return "1er"
        else:
            return str(j)
    
    a1, m1, j1 = str(d1.year), month_str[d1.month], day_str(d1.day)
    a2, m2, j2 = str(d2.year), month_str[d2.month], day_str(d2.day)
    
    if a1==a2 and m1==m2:
        return  j1+" au "+j2+" "+m2+" "+a2
    elif a1==a2 and m1!=m2:    
        return  j1+" "+m1+" au "+j2+" "+m2+" "+a2 
    else:
        return  j1+" "+m1+" "+a1+" au "+j2+" "+m2+" "+a2


def arrondie_str(x):
    corps, decimales = str(x).split('.')
    return corps+','+decimales[:2]

def moyennes_annuelles(data, date_depart, periode=timedelta(7)):
    date1 = date_depart-periode
    # 1 an avant la date d'analyse:
    date4 = date_depart-52*timedelta(7)
    date3 = date4-periode
    # 2 ans avant la date d'analyse
    date6 = date_depart-104*timedelta(7)
    date5 = date6-periode
    
    moy12 = data[(data.index>date1) & (data.index<=date_depart)].mean()
    moy34 = data[(data.index>date3) & (data.index<=date4)].mean()
    moy56 = data[(data.index>date5) & (data.index<=date6)].mean()
    df = pd.concat([moy56, moy34, moy12], axis=1)
    df.columns = [date6.year, date4.year, date_depart.year]
    return df.T

def variations_annuelles(data, date_depart, periode=timedelta(7)):
    date1 = date_depart-periode
    # 1 an avant la date d'analyse:
    date4 = date_depart-52*timedelta(7)
    date3 = date4-periode
    # 2 ans avant la date d'analyse
    date6 = date_depart-104*timedelta(7)
    date5 = date6-periode
    
    moy12 = data[(data.index>date1) & (data.index<=date_depart)].mean()
    moy34 = data[(data.index>date3) & (data.index<=date4)].mean()
    moy56 = data[(data.index>date5) & (data.index<=date6)].mean()
    df = pd.concat([(moy12-moy56)/moy56*100,
                    (moy12-moy34)/moy34*100], axis=1)
    df.columns = [str(date_depart.year) +" vs "+str(date6.year),
                  str(date_depart.year) +" vs "+str(date4.year)]
    return df.T


### III - CALCULS
def variation(x, delta=timedelta(days=7)):
    t2 = max(x.index)
    t1 = t2-delta
    return (x[t2]-x[t1])/x[t1]


def variations(data, date1, date2, delta=4*timedelta(7)):
    # Variations en pourcentage
    dt = data.index[-1] - data.index[-2]
    var = 100*(data-data.shift(round(delta/dt)))/data.shift(round(delta/dt))

    # Variations pendant delta, pour toutes les dates entre date1 et date2 
    date_min = max(min(data.index), date1-delta)
    date_max = min(max(data.index), date2)
    var = var[(var.index>=date_min) & (var.index<=date_max)]
    
    # double index avec la date de début et de fin 
    #dates_1, dates_2 = var.index-delta, var.index
    #dates_1.name, dates_2.name = "début", "fin"
    #var.index = [dates_1, dates_2]

    return var


def tops3(data, date1, date2):

    def tops(data, date1, date2):
        data = data[(data.index>=date1) & (data.index<=date2)]
        tops = data.mean().sort_values(ascending=False)
        return tops
    
    var = variations(data, date1, date2, delta=4*timedelta(7))

    tops_volume    = tops(data, date1, date2)
    tops_variation = tops(var , date1, date2)
    tops_potentiel = (tops_variation*tops_volume).sort_values(ascending=False)

    tops3 = pd.DataFrame({
        "top volume"      : list(tops_volume.head(3).index),
        "top progression" : list(tops_variation.head(3).index),
        "top potentiel"   : list(tops_potentiel.head(3).index)}).T

    # tops3 = tops3.applymap(lambda x: nom_pays(x)+"("+x+")")
    tops3.columns = ["1er", "2ème", "3ème"]
    
    return tops3


def tops_pays(recapitualif_x_semaines, fichier, str_top_semaine):
    """ Fonction retournant un tableau du top 3 des pays ayant le plus gros
    Volume, d'un top 3 des pays ayant le plus haut top de progression ainsi
    qu'un top 3 des pays ayant le plus de potentiel
    Exemple:
        top Volume       top Progression        Top Potentiel
    0  'FR', 'BE', 'NL'  'CH', 'IT', 'NL'  'IT', 'CH', 'NL'
    
    recapitualif_x_semaines: dataframe du classement sur x semaines
    exemple: recap_desc_2s 
    et le top_semaine étant le nom de la colonne 
    en string
    exemple: "TOP 2 SEMAINES"
    """
    top = {"top Volume": [], "top Progression": [], "Top Potentiel": []} 
    
    recapitualif_x_semaines = recapitualif_x_semaines.sort_index()
    recapitualif_x_semaines.fillna(0, inplace=True)
    variation = (variations(fichier, 1).T).sort_index()
    variation.fillna(0, inplace=True)
    concat_tableau = pd.concat([variation, recapitualif_x_semaines], axis=1)
    top_volume = recapitualif_x_semaines.head(3).index.to_list()
    top_progression = variation.sort_values(by=list(variation.columns), 
                                            ascending=False).head(3).index.to_list()
    
    concat_tableau["potentiel"] = concat_tableau[list(concat_tableau.columns)[0]]*concat_tableau[str_top_semaine]
    top_potentiel = list(concat_tableau.sort_values(by=["potentiel"]).head(3).index)
    
    def nettoyage_str(x):
        """ Fonction qui permet de remplacer les "[" ainsi que les "]"
        pour avoir un tableau identique à celui du pdf du client
        """
        x = str(x)
        if "[" and "]" in x:
            x = x.replace("[", "").replace("]", "")
        return x

    top["top Volume"].append(top_volume)
    top["top Progression"].append(top_progression)
    top["Top Potentiel"].append(top_potentiel)
    colonnes = list(top.keys())
    top_3_pays = pd.DataFrame(top, columns=colonnes)
    
    for nom in colonnes:
        top_3_pays[nom] = top_3_pays[nom].apply(nettoyage_str)    
    
    return top_3_pays

### IV - GRAPHQUES

def graph_barres(data, nom_x, nom_y, nom_z, formate_date=True):
    # Mise en forme des données (data) pour pouvoir utiliser seaborne, dans un 
    # tableau à trois colonnes (data_graph). La première est le temps, sous 
    # forme de date, la deuxième est les valeurs (volumes, variations, etc...),
    # la troixième les catégories (pays, région, etc..).
    # Les légendes des axes du dessin sont:
    # légende des catégories -> nom_x
    # légende des valeurs    -> nom_y
    # légende du temps       -> nom_z
    data_graph = pd.DataFrame()
    for pays in list(data.columns):
        df = pd.DataFrame({nom_z: data[pays].index, nom_y: data[pays], nom_x: pays})
        data_graph = data_graph.append(df, ignore_index=True)

    # Lorsque les valeurs sont des volumes, les dates représentent des 
    # semaines. Elles sont mises sous un format plus lisible.
    # Lorsque les valeurs sont des variations, les dates représentent le début
    # de la première semaine de variation 
    if formate_date:
        dt = timedelta(days=6) # temps entre le début et la fin de la semaine 
        data_graph[nom_z] = data_graph[nom_z].apply(lambda t: duree_str(t, t+dt))

    # Les volumes sont ensuite représentés à l'aide de barres.
    # Différentes palettes de couleurs ont été testées:
    # YlGnBu RdBu OrRd PRGn Spectral YlOrBr
    fig, ax = plt.subplots(figsize=(10,6), dpi=250)
    sns.barplot(x=nom_x, y=nom_y, hue=nom_z, data=data_graph,
                palette=sns.color_palette("YlGnBu")[-min(len(data),6):])

    # Les différentes semaines sont données en légende
    ax.legend(loc='lower center', bbox_to_anchor=(0.5, 1.01),
              fancybox=True, shadow=False, ncol=3)

    # Les volumes sont écrits en bleu en haut d'une barre, lorsque la valeur
    # est positive et en bas d'une barre lorsque la valeur est négative.
    for p in ax.patches:
        text = " "+format(p.get_height(), '.1f')+" "
        if "%" in nom_y: text+="% "
        x = p.get_x() + p.get_width() / 2.
        y = p.get_height()
        if y >= 0:
            ax.annotate(text, (x,y), ha='center', va='bottom', size=8, 
                        color='blue', xytext=(0,1), textcoords='offset points',
                        rotation=90)
        else:
            ax.annotate(text, (x,y), ha='center', va='top', size=8, 
                        color='red', xytext=(0,1), textcoords='offset points',
                        rotation=90)
    # Des limites un peu plus larges sont fixées en ordonnées afin d'être 
    # certain que les écritures précédentes ne dépassent du cadre
    ymin, ymax = min(data_graph[nom_y]), max(data_graph[nom_y])
    try:
        ax.set_ylim([(ymin-0.2*(ymax-ymin) if ymin < 0 else 0),
                     (ymax+0.2*(ymax-ymin) if ymax > 0 else 0)])
    except:
        pass

    plt.xticks(rotation=45)

    return fig

@st.cache()
def prevision_prophet(data,pays,nb_semaines = 4):
    
    date_fin = data.index[-1]
    data_index = data[pays].index.name
    #if lissage:
    data_cast = data.reset_index()[[data_index,'lisse']]
    #else:
        #data_cast = data.reset_index()[[data_index,pays]]

    data_cast.columns = ['ds','y']
    
    
    m = Prophet(seasonality_mode='additive',
                daily_seasonality=False,
                weekly_seasonality=False,
                yearly_seasonality=True,
                growth='linear',
                changepoint_prior_scale = 0.1, # Increasing it will make the trend more flexible
                seasonality_prior_scale = 1,
                changepoint_range=0.85)
    
    m = m.fit(data_cast)
    
    future = m.make_future_dataframe(nb_semaines,freq='W')
    
    future['cap']=0
    future['floor'] = None
    
    forecast = m.predict(future)
    
    result = forecast[['ds', 'yhat']].copy()
    result.columns = [data_index, pays]
    result.set_index(data_index,inplace = True)
    result[result < 0] = 0
    result.index = result.index.map(lambda x: x.date())
    #st.write(result[pays][-nb_semaines])
    #st.write(data_cast['y'].iloc[-1])
    try :
        last_val = data_cast['y'].iloc[-1]
        first_pred = result[pays][-nb_semaines]
        
        gap = abs((first_pred - last_val) / last_val)*100
        if gap > 5:
            # cas trop grand :
            if first_pred > last_val:
                delta = last_val + 0.05*last_val - first_pred
            else : 
                delta = last_val - 0.05*last_val - first_pred
            
            result[pays] = result[pays] + delta
            #st.write(delta)

                
    except :
        pass
    
    return result

def graph_3_ans(data, pays, lissage=False, prevision = True, nb_semaines = 0):
    """Lissage avec le filtre de Savitzky-Golay . Il utilise les moindres 
    carrés pour régresser une petite fenêtre de vos données sur un polynôme, 
    puis utilise le polynôme pour estimer le point situé au centre de la 
    fenêtre. Enfin, la fenêtre est décalée d’un point de données et le 
    processus se répète. Cela continue jusqu'à ce que chaque point ait été 
    ajusté de manière optimale par rapport à ses voisins. Cela fonctionne très 
    bien même avec des échantillons bruyants provenant de sources non 
    périodiques et non linéaires."""
    
    fig, ax = plt.subplots(figsize=(10,6), dpi=250)
    a = max(data.index).year
    j1 = data[data.index >= datetime(a, 1, 1).date()].index[0]
    jlast = data[data.index >= datetime(a, 1, 1).date()].index[-1]
    #st.write(data[pays])

    # Si "prevision" est cochée on fait la prédiction 
    
    if lissage and not prevision :
        data['lisse'] = savgol_filter(data[pays].values, 9, 3, mode='mirror')
    
    if prevision == True : 
        data['lisse'] = savgol_filter(data[pays].values, 9, 3, mode='mirror')
        data_predict = prevision_prophet(data,pays,nb_semaines = nb_semaines)
        annee_fin = data_predict.index[-1].year
        
        data_predict = data_predict[data_predict.index >= jlast]
        if lissage:
            data_predict.iloc[0] = data['lisse'].iloc[-1]
        else : 
            data_predict.iloc[0] = data[pays].iloc[-1]
        #st.dataframe(data_predict)
        taille_predict = len(data_predict)
 
        for i in range(annee_fin - jlast.year + 1):
            
            date1, date2 = datetime(a+i, 1, 1).date(), datetime(a+i, 12, 31).date()
            data_ = data_predict[(data_predict.index>=date1) & (data_predict.index<=date2)]
            if i == 0:
                label_ = 'Prévision lisse'
                dates = [jlast+ d*timedelta(days=7) for d in range(len(data_))]
            else:
                label_ = None
                label_lisse = None
                dates = [j1+ d*timedelta(days=7) for d in range(len(data_))]
                
            y = data_[pays]
            
            ax.plot(dates, y, 'o-', color='red', label=label_)
    
    
            
    for i in range(4):
        date1, date2 = datetime(a-i, 1, 1).date(), datetime(a-i, 12, 31).date()
        data_ = data[(data.index>date1) & (data.index<=date2)]
        
        dates = [j1+int((date-date1).days/7.)*timedelta(days=7) for date in data_.index]
        
        ligne  = ('o--' if i==0 else '.-')
        ligne2 = ('o:'  if i==0 else '.:')
        #c = sns.color_palette("YlGnBu")[-i*2-1]
        c = sns.color_palette("deep")[-i]
        y = data_[pays].values
        # épaisseur différente pour 2019
        width=1.5
        if i == 3 :
            width=2.4
            
        if lissage:
            ylis = data_['lisse'].values
            ax.plot(dates, ylis, ligne, color=c, label=str(a-i)+u" lissé",linewidth=width)
            ax.plot(dates, y, ligne2, color=c, label=str(a-i), alpha=0.3,linewidth=width)
        else:
            ax.plot(dates, y, ligne, color=c, label=str(a-i),linewidth=width)
    
    # Les différentes semaines sont données en légende
    ax.legend(fancybox=True, shadow=False, ncol=1)
    
    # Des limites pour que l'échelle ne change pas entre le lissage et 
    # l'abscence de lissage 
    ax.set_ylim(0, 1.1*data[pays].max())
    
    ax.set_ylabel("Indice Google Trends – Base 100")
    ax.set_title(pays)
    
    plt.xticks([datetime(a, m+1, 1).date() for m in range(12)], 
           ['Janvier', 'Février', 'Mars', 'Avril', 'Mai', 'Juin',
            'Juillet', 'Août', 'Septembre', 'Octobre', 'Novembre', 'Décembre'],
           rotation=45) 
        
    return fig


### V - GENERATION D'UN RAPPORT
def rapport_pdf():
    pass


### VI - INTERFACES WEB

def entete():
    txt = u"""Bienvenue à l’observatoire digital des destinations françaises
et européennes de XXXXXX – powered by BC.
L’observatoire digital de XXXXXX mesure, par quinzaine, par mois 
et par trimestre, les niveaux d’intérêts d’un marché
dans Google Trends (rubrique « travel  ») d’une sélection de mots clés
génériques et des destinations touristiques françaises par espaces
(littoral, outre-mer, urbain, campagne et montagne) et propose des les
comparer à la concurrence en Europe.
    """
    cols = st.columns(2) # number of columns in each row! = 2
    cols[1].image("logo_Baudy_Co.png", use_column_width=True) 
    #cols[1].image("https://nicolasbaudy.files.wordpress.com/2020/02/cropped-logo-new-2.png")
    st.title("Observatoire Digital des Destinations")
    st.text(txt)


def introduction():
    txt_1 = """
Mesure des "termes de recherches touristiques" (groupes de termes 
correspondants à un même concept) non accolés aux destinations pour mesurer
si l'intérêt pour le terme analysé redémarre, indépendamment du lieu
(ex: "hôtel" et non "hôtel à Lyon") rubrique "travel" de Google Trends.
Périodicité d'analyse: Quinzaine (2 semaines), mensuelle (4 semaines) et
trimestrielle (12 semaines)
Marchés analysés: hôtel, résidence de tourisme, camping, chambre d'hôte, 
voyage, tout inclus, week-end, croisière, billet d'avion, billet de train,
Paris et Disneyland Paris (toutes catégories)
"""

    st.title("Introduction")
    st.header("1- Analyse des mots clé génériques par pays")
    st.text(txt_1)
    
    txt_2 = """
Mesure des destinations françaises et européennes d'un panel donné
(destinations domestiques écartées sauf pour la France)
Périodicité d'analyse: Quinzaine (2 semaines), mensuelle (4 semaines) et
trimestrielle (12 semaines)
Marchés analysés: outre-mer, campagne, littoral, urbain et montagne France et
Europe - sauf outre-mer (monde)
"""
    st.header("2- Analyse toutes destinations Françaises et Européennes par pays")
    st.text(txt_2)
    
    txt_3 = """
Le lissage est une technique qui consiste à réduire les irrégularités d'une
courbe. Cette dernière est utilisée en traitement du signal pour atténuer
ce qui peut être considéré comme une perturbation ou un bruit de mesure.
"""
    st.header("3- Lissage")
    st.text(txt_3)
    


def visualisation_tops(data):
    date_1, date_2 = max(data.index) - 4*timedelta(7), max(data.index)
    txt = f"""
Synthèse des classements des 3 pays les plus dynamiques sur la période donnée, par défaut 
les 4 dernières semaines disponibles du {duree_str(date_1,date_2)}, respectivement 
pour le 'top volume', la 'top progession' et le 'top potentiel'. 
 
  - L'indicateur de 'volume' est la moyenne des volumes hebdomadaires constatés sur 
les 4 dernières semaines. Il rend compte du niveau d'activité général, tout en 
minimisant les fluctuations pouvant survenir à l'échelle hebdomadaire.
 
  - L'indicateur de 'progression' est la moyenne sur la période des variations
hebdomadaires en pourcentages. Plus il y a eu de variations hebdomadaires à la
hausse pendant 4 semaines, plus l'indicateur de progression est élévé. 

  - L'indicateur de 'potentiel' est le produit de l'indicateur de volume par 
l'indicateur de progression. Il indique les gains potentiels futurs si la tendance 
à la progression observée est conservée.
"""
    st.title("1 - Tops pays")
    st.text(txt)

    date_1, date_2 = max(data.index) - 4*timedelta(7), max(data.index)
    date1 = st.date_input("début:",value=date_1)
    date2 = st.date_input("fin:",  value=date_2)

    top3 = tops3(data, date1, date2)
    
    st.table(top3)
    ax = plt.subplot(111, frame_on=False) # no visible frame
    ax.xaxis.set_visible(False)  # hide the x axis
    ax.yaxis.set_visible(False)  # hide the y axis
    
    return top3.to_latex()


def visualisation_volumes(data):
    txt = """
Google Trends permet de mesurer, de manière relative, l’évolution des recherches 
des internautes, à partir de mots-clés (sujets ou destinations), avec un indice 100
pour la valeur la plus haute au cours de la période analysée. Le champ d’application 
est restreint au domaine du « travel » (ou catégorie « voyage »). Les résultats ne 
sont pas des valeurs absolues mais se lisent en indices.

La visualisation de cet indice au cours des dernières semaines permet de constater 
les fluctuations et les éventuelles tendances de manière empirique. L'attention est 
mise sur les 2 denières semaines puis sur les 4 dernières semaines. """
    try:
        st.title("2. Volume des tendances de recherches des deux et quatre dernières semaines ")
        st.text(txt)
    except:
        pass

    titre_googletrend = "a - Tendances de recherche des 2 dernières semaines"
    table = data.tail(2).applymap(lambda x: "{:.1f}".format(x))
    table.index = table.index.map(lambda t: duree_str(t, t+timedelta(days=6)))
    try:
        st.header(titre_googletrend)
        st.write(table)
    except:
        pass

    nom_x, nom_y, nom_z = "Pays", "Indice Google Trends – Base 100", "Semaine"
    graph_googletrends = graph_barres(data.tail(2), nom_x, nom_y, nom_z)
    try:
        st.pyplot(graph_googletrends)
    except:
        pass

    titre_tendances = "b - Tendances de recherche des 4 dernières semaines"
    table = data.tail(4).applymap(lambda x: "{:.1f}".format(x))
    table.index = table.index.map(lambda t: duree_str(t, t+timedelta(days=6)))
    graph_tendances = graph_barres(data.tail(4), nom_x, nom_y, nom_z)
    try:
        st.header(titre_tendances)
        st.write(table)
        st.pyplot(graph_tendances)
    except:
        pass
    
    resultats = {titre_tendances: graph_tendances,
                 titre_googletrend: graph_googletrends}
    return resultats 


def visualisation_variations(data):
    date = lambda i: max(data.index) + i*timedelta(7)
    semaine =lambda i: duree_str(data.index[-i], data.index[-i]+timedelta(6))
    periode = lambda i, j: "semaine du "+semaine(i)+" à la semaine du "+semaine(j)
    txt = """
D’une semaine à l’autre, les indices des tendances de recherches de Google Trends 
peuvent fluctuer. Les variations sont mesurées Semaine S vs Semaine S-1 et
Semaine S-1 vs Semaine S-2. Les variations S-1 vs S-2 sont comparées à celles
de S-2 vs S-3. """

    st.title("3 - Variations de l'indice")
    st.text(txt)

    st.header("a - Variations S/S-1 comparées à S-1/S-2")
    #st.text(f"{periode(-1,0)}")
    var = variations(data, date(-1), date(0), delta=timedelta(7)).tail(2)
    table = var.applymap(lambda x: "{:.1f}".format(x))
    #table.index = table.index.map(lambda t: duree_str(t, t+timedelta(days=6)))
    table.index = ["S-1/S-2", "S/S-1"]
    st.write(table)

    nom_x, nom_y, nom_z = "Pays", "Variation de l'indice Google Trends – %", "Semaine"
    st.pyplot(graph_barres(var, nom_x, nom_y, nom_z))

    st.header("b - Variations S-1/S-2 comparées à S-2/S-3")
    #st.text(f"{periode(-2,-1)}")
    var = variations(data, date(-2), date(-1), delta=timedelta(7)).tail(2)
    table = var.applymap(lambda x: "{:.1f}".format(x))
    #table.index = table.index.map(lambda t: duree_str(t, t+timedelta(days=6)))
    table.index = ["S-2/S-3", "S-1/S-2"]
    st.write(table)

    nom_x, nom_y, nom_z = "Pays", "Variation de l'indice Google Trends – %", "Semaine"
    st.pyplot(graph_barres(var, nom_x, nom_y, nom_z))


def interface():

    
    # def affiche_classement(destinations):
    #     """Remplace l'affichage du pays ('FR','US'...) par le classement 
    #     qu'occupe la destination, dans les noms de colonne(destinations).
    #     """
    #     place = 1
    #     nouv_colonnes = []
    #     for destination in destinations:
    #         nouv_destination = destination
    #         remplace = destination[destination.find("(")+1 : destination.find(")")]
    #         nouv_destination = destination.replace(remplace, str(place))
    #         place += 1

    #         nouv_colonnes.append(nouv_destination)
    #     return nouv_colonnes
    
    data, dossiers_source = acquisition_donnees()
    
    entete()
    
    if st.sidebar.checkbox("Présentation", value=False):
        introduction()
     
    # Bouton d'export en powerpoint, masqué pour la version client
    #if st.sidebar.button("Export ppt"):
         #export_ppt(generation_generique=True, generation_generique_par_pays=True)
        
    # Sélection du type d'analyse générale à effectuer
    types_analyse = {"Mots clés génériques": data[dossiers_source[0]],
                     "Destinations Françaises": data[dossiers_source[1]],
                     "Destinations Françaises et Européennes": data[dossiers_source[2]]}
    txt = "Types d'analyses: " 
    noms_types = list(types_analyse.keys())
    mode = st.sidebar.selectbox(txt, noms_types)

    mode = "maintenance"
    
        
    ### ANALYSE GENERIQUE
    if mode == "maintenance":
        st.markdown("**Ce service est suspendu. Merci de vous rapprocher des équipes de l’observatoire.**")
        
    elif mode == noms_types[0]:
        # Récupération des noms de tables d'analyse et construction de la 
        # liste déroulante
        noms_analyses = list(types_analyse[mode].keys())
        fichier = st.sidebar.selectbox("Quelle analyse effectuer?", noms_analyses)
        data = lecture_donnees(types_analyse[mode][fichier])
        try:
            ### 1 - LES TOPS
            if st.sidebar.checkbox("1 - Les tops") and fichier != "None":
                top3 = visualisation_tops(data)
            
            ### 2 - LES VOLUMES
            if st.sidebar.checkbox("2 - Les volumes") and fichier != "None":
                # Traçage des graphiques
                graph_volumes = visualisation_volumes(data)
                    
            ### 3 - LES VARIATIONS
            if st.sidebar.checkbox("3 - Les variations") and fichier != "None":
                visualisation_variations(data)
    
            # COMMENTAIRE à inclure dans le rapport
            #if st.checkbox("Voulez-vous mettre un commentaire?"):
            #    commentaire_1 = st.text_area("Commentaire", "")
    
        except:
            pass

    ### ANALYSE PAR PAYS
    else:
        tous_pays = list(types_analyse[mode].keys())
        pays_choisi = st.sidebar.selectbox("Quel pays?", tous_pays)
        detail_analyse = list(types_analyse[mode][pays_choisi].keys())
        detail_analyse = changement_OM(detail_analyse)
        analyse_pays = st.sidebar.selectbox("Quelle analyse effectuer?",
                                           detail_analyse)
        analyse_pays = find_key(analyse_pays)
        data = lecture_donnees(types_analyse[mode][pays_choisi][analyse_pays])
        
        try:
            # Date d'analyse
            txt = "Date d'analyse"
            date2 = st.sidebar.date_input(txt,value=max(data.index))

            # Moyennes des volumes sur 2, 4 et 12 semaines,
            # triés par ordre décroissant
            moyennes = {}
            for i in [2, 4, 12]:
                date1 = date2 - i*timedelta(7)
                calcul_moy = data[(data.index>date1) & (data.index<=date2)].mean()
                calcul_moy = calcul_moy.sort_values(ascending=False)
                calcul_moy.name = "TOP " + str(i) + " SEMAINES"
                moyennes[i] = calcul_moy
                
            ### 1 - LES TOPS
            if st.sidebar.checkbox("1- Les tops") and analyse_pays != "None":
                st.title("1 - Les tops tendances de recherche - Base : indice 100")
                txt = f"""
Les valeurs moyennes des tendances de recherche de Google Trends sont classées,
sur des périodes, de respectivement:
    - 2 semaines, du {duree_str(date2- 2*timedelta(7), date2)}
    - 4 semaines, du {duree_str(date2- 4*timedelta(7), date2)}
    - 12 semaines, du {duree_str(date2-12*timedelta(7), date2)}"""
                st.text(txt)

                st.header("a - Le top 6")
                cols, k = st.columns(3), 0
                for i, k in zip([2, 4, 12],[0,1,2]):
                    cols[k].table(moyennes[i].apply(arrondie_str).head(6))

                if st.checkbox("afficher les valeurs suivantes..."):
                    st.header("b - Les valeurs suivantes")
                    cols, k = st.columns(3), 0
                    for i, k in zip([2, 4, 12],[0,1,2]):
                        cols[k].table(moyennes[i].apply(arrondie_str).iloc[7:])

                # # COMMENTAIRE à inclure dans le rapport
                # if st.checkbox("Voulez vous mettre un commentaire ?"):
                #     commentaire_2 = st.text_area("Commentaire", "")
           
            ### 2 - LES VOLUMES
            if st.sidebar.checkbox("2 - Les volumes des 4 dernières années"):
                st.title("2 - Les volumes des 4 dernières années")
                # Le type d'analyse peut être traduit en un nombre de semaines
                classements = {'2 semaines': 2,
                                '4 semaines': 4,
                                '12 semaines': 12}
                types_classement = list(classements.keys())
                choix_vol = st.sidebar.radio("Classement: ", types_classement)
                lissage = st.sidebar.checkbox("Lissage")
                
                # Le choix est donné pour n'afficher que les destinations
                # voulues parmi toutes celles disponibles
                nb_semaines_vol = classements[choix_vol]
                choix_destinations = {}
                correspond_vol = {}
                max_colonnes = 5
                colonnes_volume = st.columns(max_colonnes)
                index = 0  
                place = 1
                for destination in moyennes[nb_semaines_vol].index:
                    nom_classe = destination[:destination.find("(")]
                    nom_classe += " (" + str(place) + ")"
                    correspond_vol[nom_classe] = destination
                    choix_destinations[destination] = colonnes_volume[index].checkbox(nom_classe)
                    index += 1
                    place += 1
                    if index == max_colonnes:
                        index = 0
                        
                # Prévision
                
                prevision = st.sidebar.checkbox('Afficher les prévisions')
                nb_semaines = 0
                if prevision:
                    nb_semaines = st.sidebar.number_input(
                        "Horizon de prévision (en semaines) ",
                        min_value = 1, max_value = 16,value=4)
                        
                # Seules les graphiques des destinations choisies sont affichés
                # choix_vol = st.multiselect("Choisissez les données à analyser:", correspond_vol)
                
                for zone in choix_destinations:
                    if choix_destinations[zone] == True:
                        st.pyplot(graph_3_ans(data, zone, lissage,prevision,nb_semaines))
    
                # for zone in correspond_vol:
                #     if zone in choix_vol:
                #         st.pyplot(graph_3_ans(data, correspond_vol[zone], lissage))
                        

            ### 3 - LES VARIATIONS
            titre_variation = "3 - Les variations des recherches d'une année sur l'autre"
            if st.sidebar.checkbox(titre_variation):
                txt = f"""
Les indices de Google Trends, moyennés au choix sur 2, 4 ou 12 dernières semaines
précédant la date d'analyse, sont comparées aux indices sur les mêmes périodes
des années précedentes."""
            
                st.title(titre_variation)
                st.text(txt)
                classements = {'2 semaines': 2,
                               '4 semaines': 4,
                               '12 semaines': 12}
                types_classement = list(classements.keys())
                choix_var = st.sidebar.radio("Moyennes sur: ", types_classement)
                nb_semaines_var = classements[choix_var]
                date1 = date2 - nb_semaines_var * timedelta(7)
                
                # Les pays dont on souhaite visualiser les variations
                # peuvent être sélectionnés et sont tous affichés ensemble
                # dans deux graphiques
                max_colonnes_var = 5
                colonnes_variation = st.columns(5)
                index_var = 0
                zones = []
                correspond_var = {}
                place = 1
                choix_variations = {}
                for destination in moyennes[nb_semaines_var].index:
                    nom_classe = destination[:destination.find("(")]
                    nom_classe += " (" + str(place) + ") "
                    correspond_var[nom_classe] = destination
                    choix_variations[destination] = colonnes_variation[index_var].checkbox(nom_classe)
                    index_var += 1
                    place += 1
                    if index_var == max_colonnes_var:
                        index_var = 0
                                        
                for choix in choix_variations:
                    if choix_variations[choix] == True:
                        zones.append(choix)
                        
                moy = moyennes_annuelles(data[zones], date2,
                                         nb_semaines_var * timedelta(7))
                var = variations_annuelles(data[zones], date2,
                                           nb_semaines_var * timedelta(7))
                
                # Graphique des moyennes comparées
                """On décale d'une semaine les deux dates,pour l'affichage"""
                """On décale d'une semaine les deux dates, pour l'affichage"""
                date1bis = date1 + timedelta(7)
                date2bis = date2 + timedelta(7-1)
                titre_var = "a) Valeurs du " + duree_str(date1bis, date2bis)
                titre_var += " comparées aux années précédentes."
                st.header(titre_var)                
                st.table(moy.T.applymap(lambda x: "{:.1f}".format(x)))
                nom_x, nom_z = u"Régions", "Annees"
                nom_y = "Moyennes de l'indice Google Trends"
                st.pyplot(graph_barres(moy, nom_x, nom_y, nom_z,
                                       formate_date=False))
                
                # Graphiques des variations
                st.header("b) Variations en %")
                st.table(var.T.applymap(lambda x: "{:.1f}".format(x)))
                nom_y = "Variation des moyennes de l'indice Google Trends - %"
                st.pyplot(graph_barres(var, nom_x, nom_y, nom_z,
                                       formate_date=False))

        except:
            pass

### EXPORT POWERPOINT
def export_ppt(generation_generique=True, generation_generique_par_pays=True):
    def ajout_titre(page, type_analyse="", position=0,
                    titre="Indice hebdomadaire des tendances de recherches"):
        """Placement du titre de page. Selon si on on précise un type d'analyse
        ou pas, il sera composé avec une rubique ou non.
        Le titre prend par défaut la première position, mais il est également 
        possible de le placer en dessous d'un autre contenu, avec une position 
        plus élevée."""
        shape = page.shapes[position]
        text_frame = shape.text_frame
        p = text_frame.paragraphs[0]
        p.margin_left = 0
        run = p.add_run()
        run.text = titre
        if type_analyse != "":
            run.text += ' . Rubrique ' + type_analyse
        font = run.font
        font.name = 'Calibri'
        font.size = Pt(18)
        font.bold = True
        font.italic = None
        font.color.rgb = RGBColor(0x11, 0x55, 0xCC)
        
    def table_ppt(page, data, place_page=0, pos_y=1.5, hauteur=3.5):
        """Création d'une table"""
        nb_colonnes, nb_lignes = data.shape[1], data.shape[0]
        taille = nb_colonnes * 1.5
        x, y, cx, cy = Inches(0.5), Inches(pos_y), Inches(taille), Inches(hauteur)
        shape = page.shapes.add_table(nb_lignes+1, nb_colonnes, x, y, cx, cy)       
        table = shape.table
        index_col = 0
        for nom_colonne in data.columns:
            table.cell(0, index_col).text = nom_colonne
            index_ligne = 1
            for valeur in data[nom_colonne].tolist():
                table.cell(index_ligne, index_col).text = str(valeur)
                index_ligne += 1
            index_col += 1
            
        def iter_cells(table):
            for row in table.rows:
                for cell in row.cells:
                    yield cell
        
        for cell in iter_cells(table):
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.size = Pt(10)
                    
    def calcul_tops(analyses, nb_semaines):
        """Création d'un tableau des tops volumes, progression et potentiel.
        Paramètres:
            analyses: Dictionnaire -> récapitulatif des DataFrames pour tous
                        les types d'analyses: chiffres par pays ou par destination.
            nb_semaines: Entier -> nombre de semaines depuis la dernière date
                            pour effectuer le calcul."""
        colonnes = ["", "Top Volume", "Top Progression", "Top Potentiel"]
        tops = pd.DataFrame(columns=colonnes)
        for type_analyse in analyses:
            analyse = lecture_donnees(analyses[type_analyse])
            date_2 = max(analyse.index)
            date_1 = date_2 - nb_semaines*timedelta(7)
            top = tops3(analyse, date_1, date_2)
            volume = ",".join(top.loc['top volume'])
            progression = ",".join(top.loc['top progression'])
            potentiel = ",".join(top.loc['top potentiel'])
            tops.loc[len(tops.index)] = [type_analyse,
                                         volume,
                                         progression,
                                         potentiel]
        return tops

       
    data, dossiers_source = acquisition_donnees()
    
    
    # Générique
    if generation_generique:
        print("Génération du PowerPoint Générique")
        presente = Presentation()
        page_titre = presente.slide_layouts[1]
        slide = presente.slides.add_slide(page_titre)
        slide.shapes.add_picture("logo_Baudy_Co.png",
                                  Inches(4.5), Inches(3),
                                  width = Inches(5))
        titre = slide.shapes.title
        titre.text = u"""Observatoire digital des destinations
        Analyse Générique"""
        
        
        # Page des priorités d'action
        page_priorite = presente.slides.add_slide(page_titre)
        
        date_1, date_2 = "", ""
        colonnes = ["", "Top Volume", "Top Progression", "Top Potentiel"]
        top_quinzaine = pd.DataFrame(columns=colonnes)
        for type_analyse in data['generiques']:
            analyse = lecture_donnees(data['generiques'][type_analyse])
            date_2 = max(analyse.index)
            date_1 = date_2 - 2*timedelta(7)
            top = tops3(analyse, date_1, date_2)
            volume = ",".join(top.loc['top volume'])
            progression = ",".join(top.loc['top progression'])
            potentiel = ",".join(top.loc['top potentiel'])
            top_quinzaine.loc[len(top_quinzaine.index)] = [type_analyse, volume,
                                                            progression, potentiel]
        top_priorite = top_quinzaine[["", "Top Progression"]]
        top_priorite.columns = ["", "Priorité d'action"]
                    
        titre = "La quinzaine du " + duree_str(date_1, date_2) + " en quelques mots..."
        ajout_titre(page_priorite, titre=titre, position=0)
        table_ppt(page_priorite, top_priorite, 1)
        
        # Page des tops de la quinzaine
        page_top = presente.slides.add_slide(page_titre)
                    
        titre = "La quinzaine du " + duree_str(date_1, date_2) + " en quelques mots..."
        ajout_titre(page_top, titre=titre, position=0)
        table_ppt(page_top, top_quinzaine, top_quinzaine.shape[1],
                  top_quinzaine.shape[0], 1)
        
        # Graphiques d'analyse générale
        for type_analyse in data['generiques']:
            page_analyse = presente.slides.add_slide(page_titre)
            left = top = Inches(0)
            width = Inches(10.0)
            height = Inches(0.2)
            shape = page_analyse.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, left, top, width, height
            )
            
            ajout_titre(page_analyse, type_analyse)
        
            donnees_propres = lecture_donnees(data['generiques'][type_analyse])
            graphiques = visualisation_volumes(donnees_propres)
            
            decalage = 0
            for type_graphique in graphiques:
                try:
                    graph = graphiques[type_graphique]
                    nom_graph = str(type_analyse) +" "+ str(type_graphique)+".jpg"
                    image_graph = graph.savefig(nom_graph, dpi=300)
                    place_image = page_analyse.shapes.add_picture(nom_graph,
                                                          Inches(decalage),
                                                          Inches(2),
                                                          width=Inches(5))
                    decalage += 5
                except:
                    pass
            # break # Arrêt de boucle pour test
        presente.save('Rapport analyse generique.pptx')
            
    # Génération par pays, en destinations françaises et toutes nationalités
    if generation_generique_par_pays:
        print("Génération du PowerPoint par Pays")
        for dossier in dossiers_source[1:]:
            for pays in data[dossier]:
                # Les calculs des meilleurs secteurs trimestriels, puis mensuels
                # et enfin bi-hebdomadaires, sont effectués. Trois fichiers 
                # correspondants seront produits.
                periodes = ("hebdomadaire", "mensuelle", "trimestrielle")
                nb_semaine = (2, 4, 12)
                for periodicite, nb_semaines in zip(periodes, nb_semaine):
                    presente_pays = Presentation()
                    page_titre = presente_pays.slide_layouts[1]
                    page_titre_pays = presente_pays.slides.add_slide(page_titre)
                    left = top = Inches(0)
                    width = Inches(10.0)
                    height = Inches(0.2)
                    barre = page_titre_pays.shapes.add_shape(
                                MSO_SHAPE.RECTANGLE, left, top, width, height
                            )
                    
                    titre_pays = page_titre_pays.shapes.title
                    titre_pays.text = f"""Analyse {periodicite} par Pays
                    {pays}"""
                    titre = "En quelques mots..."
                    ajout_titre(page_titre_pays, position=1, titre=titre)
                    
                    tops = calcul_tops(data[dossier][pays], 2)
                    table_ppt(page_titre_pays, tops, pos_y=2.5)
                    
                    for type_analyse in data[dossier][pays]:
                        data_pays = lecture_donnees(data[dossier][pays][type_analyse])
                        moyennes = {}
                        date2 = data_pays.index.max()
                        for i in [2, 4, 12]:
                            date1 = date2-i*timedelta(7)
                            moyennes[i] = data_pays[(data_pays.index>date1) & (data_pays.index<=date2)].mean()
                            moyennes[i] = moyennes[i].sort_values(ascending=False)
                            moyennes[i].name = "TOP "+str(i)+" SEMAINES"
                        date1 = date2 - nb_semaines*timedelta(7)
                        zones = list(moyennes[2].head(6).index)
                        moy = moyennes_annuelles(data_pays[zones], date2, nb_semaines*timedelta(7))
                        var = variations_annuelles(data_pays[zones], date2, nb_semaines*timedelta(7))
                        
                        # Graphiques en barres des moyennes et variations
                        for analyse, nom_analyse in zip((moy, var), ("Moyenne", "Variation")):
                            page_analyse = presente_pays.slides.add_slide(page_titre)
                            left = top = Inches(0)
                            width = Inches(10.0)
                            height = Inches(0.2)
                            barre = page_analyse.shapes.add_shape(
                                        MSO_SHAPE.RECTANGLE, left, top, width, height
                                    )
                            ajout_titre(page_analyse, position=0, type_analyse=type_analyse)
                            nom_x, nom_z = u"Régions", "Annees"
                            nom_y = nom_analyse + " de l'indice Google Trends"
                            graph = graph_barres(analyse, nom_x, nom_y, nom_z,
                                                    formate_date=False)
                            nom_graph =  " ".join([nom_analyse,periodicite,
                                                    str(type_analyse),str(pays)])+".jpg"
                            image_graph = graph.savefig(nom_graph, dpi=250,
                                                        bbox_inches="tight")
                            plt.clf()
                            plt.cla()
                            plt.close('all')
                            del graph
                            page_analyse.shapes.add_picture(nom_graph,
                                            Inches(1),
                                            Inches(1.3),
                                            width=Inches(8))
                            
                        # Graphique en ligne
                        page_analyse = presente_pays.slides.add_slide(page_titre)
                        left = top = Inches(0)
                        width = Inches(10.0)
                        height = Inches(0.2)
                        barre = page_analyse.shapes.add_shape(
                                    MSO_SHAPE.RECTANGLE, left, top, width, height
                                )
                        ajout_titre(page_analyse, position=0, type_analyse=type_analyse)
                        cale_gauche = 1.8
                        cale_haut = 1.3
                        decalage_x = cale_gauche
                        decalage_y = cale_haut
                        for colonne in moy.columns:
                            graph = graph_3_ans(data_pays, colonne)
                            nom_graph = " ".join(("Evolution",periodicite, pays,
                                                  str(colonne))) + ".jpg"
                            image_graph = graph.savefig(nom_graph, dpi=250,
                                                        bbox_inches="tight")
                            plt.clf()
                            plt.cla()
                            plt.close('all')
                            del graph
                            page_analyse.shapes.add_picture(nom_graph,
                                                Inches(decalage_x),
                                                Inches(decalage_y),
                                                width=Inches(3))
                            if decalage_x > cale_gauche:
                                decalage_x = cale_gauche
                                decalage_y += 2
                            else:
                                decalage_x += 3.3
                    presente_pays.save(f'Rapport analyse {dossier} {periodicite} {pays}.pptx')
                    # break # Arrêt de boucle pour test
                # break # Arret de boucle pour test
        

### VI - TESTS UNITAIRES
test = False

if test:
    print("lecture des données:")
    try:
        fichier = "../FR-IT-NL-GB-US-BE-CH-DE-ES_Generique-Avion-Hebdo_20210621_1048.csv"
        data = lecture_donnees(fichier)
    except:
        data = donnees_aleatoires(t0=datetime(2017, 6, 1), nb_semaines=4*53)
    print(data)
    
    print("\ntest d'écriture des noms de pays à patir des codes iso:")
    for x in ['FR', 'BE', 'IT', 'CH', 'NL', 'US', 'GB']:
        print("\tcode iso:", x, "=> nom du pays:", x)

    print("\ntest d'écriture d'une durée:")
    date1 = datetime(2021, 5,  9).date()
    date2 = datetime(2021, 5, 30).date()
    print("\tdu", date1, " au ", date2, ": ", duree_str(date1, date2))



### VII - PROGRAMME PRINCIPAL
interface()
# export_ppt(generation_generique=False)
